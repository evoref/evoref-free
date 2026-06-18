"""PPTX Writer

.pptx: ContentBlock → PowerPoint スライド
python-pptx を使用。
"""

from __future__ import annotations

import io

from backend.export._writer_base import BytesWriterBase
from backend.export.base import ContentBlock, ExportContent


def _build_pptx(content: ExportContent) -> bytes:
    """ExportContent を PPTX バイトデータに変換"""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide_layout_title = prs.slide_layouts[0]   # タイトルスライド
    slide_layout_content = prs.slide_layouts[1]  # タイトル + コンテンツ

    blocks = content.blocks
    if not blocks and content.raw_markdown:
        from backend.export.content_converter import ContentConverter
        blocks = ContentConverter().convert(content.raw_markdown)

    # ブロックをスライド単位にグループ化
    # heading (level 1-2) でスライド分割
    slides_data: list[tuple[str, list[ContentBlock]]] = []
    current_title = content.title or ""
    current_blocks: list[ContentBlock] = []

    for block in blocks:
        if block.type == "heading" and block.level <= 2:
            if current_title or current_blocks:
                slides_data.append((current_title, current_blocks))
            current_title = block.content
            current_blocks = []
        else:
            current_blocks.append(block)

    if current_title or current_blocks:
        slides_data.append((current_title, current_blocks))

    # スライドが空の場合はタイトルスライドだけ作成
    if not slides_data:
        slide = prs.slides.add_slide(slide_layout_title)
        slide.shapes.title.text = content.title or "Untitled"
    else:
        for title, slide_blocks in slides_data:
            slide = prs.slides.add_slide(slide_layout_content)
            slide.shapes.title.text = title

            # コンテンツプレースホルダーにテキストを追加
            body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            if body is None:
                continue

            tf = body.text_frame
            tf.clear()

            for i, block in enumerate(slide_blocks):
                if block.type == "paragraph":
                    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    para.text = block.content

                elif block.type == "list":
                    for item in block.items:
                        para = tf.paragraphs[0] if (i == 0 and not tf.paragraphs[0].text) else tf.add_paragraph()
                        para.text = item
                        para.level = 1

                elif block.type == "code":
                    para = tf.paragraphs[0] if (i == 0 and not tf.paragraphs[0].text) else tf.add_paragraph()
                    para.text = block.content
                    for run in para.runs:
                        run.font.name = "Consolas"
                        run.font.size = Pt(10)

                elif block.type == "table":
                    # テーブルはプレースホルダー外に配置
                    if block.rows:
                        rows_count = len(block.rows)
                        cols_count = len(block.rows[0])
                        left = Inches(0.5)
                        top = Inches(3.5)
                        width = Inches(9.0)
                        height = Inches(0.3 * rows_count)
                        table = slide.shapes.add_table(
                            rows_count, cols_count, left, top, width, height,
                        ).table
                        for r_idx, row_data in enumerate(block.rows):
                            for c_idx, cell_text in enumerate(row_data):
                                table.cell(r_idx, c_idx).text = cell_text

                elif block.type == "quote":
                    para = tf.paragraphs[0] if (i == 0 and not tf.paragraphs[0].text) else tf.add_paragraph()
                    para.text = f'"{block.content}"'
                    para.font.italic = True

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class PptxWriter(BytesWriterBase):
    """PPTX ファイル Writer"""

    @property
    def extensions(self) -> frozenset[str]:
        return frozenset({".pptx"})

    @property
    def requires(self) -> list[str]:
        return ["python-pptx"]

    def _render_bytes(self, content: ExportContent, ext: str) -> bytes:
        return _build_pptx(content)

    def is_available(self) -> bool:
        try:
            from pptx import Presentation  # noqa: F401
            return True
        except ImportError:
            return False
