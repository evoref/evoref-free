"""PPTX Extractor

.pptx ファイルからスライドのテキストを抽出する。
python-pptx を使用。
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, BinaryIO

from backend.extraction._binary_source_base import BinarySourceExtractorBase
from backend.extraction.base import ExtractionError


class PptxExtractor(BinarySourceExtractorBase):
    """PPTX ファイルからテキストを抽出"""

    @property
    def extensions(self) -> frozenset[str]:
        return frozenset({".pptx"})

    @property
    def error_code(self) -> str:
        return "pptx_error"

    @property
    def requires(self) -> list[str]:
        return ["python-pptx"]

    def is_available(self) -> bool:
        try:
            from pptx import Presentation  # noqa: F401
            return True
        except ImportError:
            return False

    def _extract_from_source(
        self,
        source: Path | BinaryIO,
        source_name: str,  # noqa: ARG002
    ) -> tuple[str, dict[str, Any]]:
        Presentation = self._import_pptx()
        # python-pptx は str / file-like を受け付ける。
        # 既存挙動 (str(path)) を維持するため Path のときのみ str に変換する
        arg = str(source) if isinstance(source, Path) else source
        prs = Presentation(arg)
        text, slide_count = self._extract_slides(prs)
        return text, {"slides": slide_count}

    @staticmethod
    def _import_pptx():
        try:
            from pptx import Presentation
            return Presentation
        except ImportError:
            raise ExtractionError(
                "missing_library",
                "python-pptx is required for PPTX extraction: pip install python-pptx",
            )

    @staticmethod
    def _extract_slides(prs) -> tuple[str, int]:
        """全スライドからテキストを抽出"""
        texts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text)
            if slide_texts:
                texts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
        return "\n\n".join(texts), len(prs.slides)
